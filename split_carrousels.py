
"""
SPLIT_CARROUSELS.PY — Carousel Engine Pro
Genere des carrousels PPTX haute qualite a partir de texte structure.
Les fichiers sont exportes en PPTX dans pptx_exports/ pour upload LinkedIn.
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# CONFIGURATION
# ============================================================
QUEUE_FILE = "content_queue.json"
PPTX_DIR = "pptx_exports"

# Design - Couleurs de marque
BRAND_COLORS = {
    "primary": RGBColor(0x1B, 0x2A, 0x4A),
    "secondary": RGBColor(0x2E, 0x86, 0xAB),
    "accent": RGBColor(0xE8, 0x4D, 0x2E),
    "background": RGBColor(0xF8, 0xF9, 0xFA),
    "text_dark": RGBColor(0x1A, 0x1A, 0x2E),
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),
    "highlight": RGBColor(0xFF, 0xD7, 0x00),
    "grey_light": RGBColor(0xE0, 0xE0, 0xE0),
    "grey_text": RGBColor(0x4A, 0x4A, 0x4A),
    "grey_subtle": RGBColor(0xAA, 0xAA, 0xAA),
}

# Dimensions LinkedIn carousel (ratio 4:5 = 1080x1350px)
SLIDE_WIDTH = Inches(7.5)
SLIDE_HEIGHT = Inches(9.375)


# ============================================================
# PARSING DU CONTENU CARROUSEL
# ============================================================
def parse_carousel_content(content):
    """
    Parse le contenu en slides individuelles.
    Formats acceptes:
    - SLIDE 1: [contenu]
    - Separation par double saut de ligne
    - Numerotation 1. 2. 3.
    """
    slides = []

    # Format "SLIDE X:"
    if "SLIDE" in content.upper():
        parts = content.upper().split("SLIDE")
        # Re-split sur le contenu original en gardant la casse
        import re
        slide_pattern = re.compile(r'SLIDE\s*\d+\s*:', re.IGNORECASE)
        raw_parts = slide_pattern.split(content)

        for part in raw_parts:
            cleaned = part.strip()
            if cleaned and len(cleaned) > 5:
                slides.append(cleaned)

    # Fallback: double saut de ligne
    if not slides:
        parts = content.split("\n\n")
        for part in parts:
            cleaned = part.strip()
            if cleaned and len(cleaned) > 5:
                slides.append(cleaned)

    # Fallback 2: numerotation
    if len(slides) < 3:
        import re
        numbered = re.split(r'\n\d+[\.\)]\s', content)
        if len(numbered) >= 3:
            slides = [p.strip() for p in numbered if p.strip()]

    return slides


def structure_slides(raw_slides):
    """Structure les slides avec type, titre et contenu."""
    structured = []

    for i, slide_text in enumerate(raw_slides):
        lines = slide_text.strip().split("\n")
        lines = [l.strip() for l in lines if l.strip()]

        if i == 0:
            # Slide 1: Hook
            structured.append({
                "type": "hook",
                "title": lines[0] if lines else "",
                "subtitle": " ".join(lines[1:]) if len(lines) > 1 else "",
            })
        elif i == len(raw_slides) - 1:
            # Derniere slide: CTA
            structured.append({
                "type": "cta",
                "title": lines[0] if lines else "Et toi ?",
                "body": " ".join(lines[1:]) if len(lines) > 1 else "",
            })
        else:
            # Slides contenu
            title = lines[0] if lines else ""
            body = "\n".join(lines[1:]) if len(lines) > 1 else ""
            structured.append({
                "type": "content",
                "title": title,
                "body": body,
                "number": i,
            })

    return structured


# ============================================================
# CREATION DES SLIDES PPTX
# ============================================================
def add_background(slide, color):
    """Ajoute un fond de couleur a une slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def create_hook_slide(slide, data, pillar):
    """Cree la slide d'accroche (slide 1)."""
    add_background(slide, BRAND_COLORS["primary"])

    # Label pilier en haut
    pillar_label = "TERRAIN" if pillar == "terrain" else "ANALYSE PROCURETECH"
    label_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1), Inches(6.5), Inches(0.5)
    )
    tf = label_box.text_frame
    tf.text = pillar_label
    for para in tf.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = BRAND_COLORS["highlight"]
        para.font.bold = True

    # Titre principal
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(6.5), Inches(4)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["text_light"]
    p.alignment = PP_ALIGN.LEFT

    # Sous-titre
    subtitle = data.get("subtitle", "")
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        p2.space_before = Pt(20)

    # "Swipe" indicator
    swipe_box = slide.shapes.add_textbox(
        Inches(4.5), Inches(8.5), Inches(2.5), Inches(0.5)
    )
    tf = swipe_box.text_frame
    tf.text = "Swipe -->"
    for para in tf.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = BRAND_COLORS["grey_subtle"]
        para.alignment = PP_ALIGN.RIGHT


def create_content_slide(slide, data, pillar):
    """Cree une slide de contenu."""
    add_background(slide, BRAND_COLORS["background"])

    # Numero de slide en grand (decoration)
    num = data.get("number", 0)
    num_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(0.3), Inches(1.5), Inches(1.5)
    )
    tf = num_box.text_frame
    tf.text = f"0{num}" if num < 10 else str(num)
    for para in tf.paragraphs:
        para.font.size = Pt(48)
        para.font.color.rgb = BRAND_COLORS["grey_light"]
        para.font.bold = True
        para.alignment = PP_ALIGN.RIGHT

    # Barre de couleur
    bar_color = BRAND_COLORS["secondary"] if pillar == "analyste" else BRAND_COLORS["accent"]
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(1.5), Pt(5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()

    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(6.5), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "")
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["text_dark"]

    # Corps
    body = data.get("body", "")
    if body:
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8), Inches(6.5), Inches(4.5)
        )
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, line_text in enumerate(body.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line_text
            p.font.size = Pt(16)
            p.font.color.rgb = BRAND_COLORS["grey_text"]
            p.space_after = Pt(12)


def create_cta_slide(slide, data):
    """Cree la slide CTA finale."""
    add_background(slide, BRAND_COLORS["primary"])

    # Titre CTA
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(6.5), Inches(2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Et toi ?")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRAND_COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER

    # Corps CTA
    body = data.get("body", "")
    if body:
        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(6.5), Inches(2)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = body
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p.alignment = PP_ALIGN.CENTER

    # Actions en bas
    action_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7), Inches(6.5), Inches(1.5)
    )
    tf = action_box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "Enregistre ce post"
    p1.font.size = Pt(16)
    p1.font.color.rgb = BRAND_COLORS["highlight"]
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Partage ton experience en commentaire"
    p2.font.size = Pt(16)
    p2.font.color.rgb = BRAND_COLORS["highlight"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)


# ============================================================
# GENERATION DU FICHIER PPTX
# ============================================================
def create_carousel_pptx(structured_slides, title="carousel", pillar="terrain"):
    """Cree un fichier PPTX complet avec toutes les slides."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Layout vide (index 6)
    blank_layout = prs.slide_layouts[6]

    for slide_data in structured_slides:
        slide = prs.slides.add_slide(blank_layout)
        slide_type = slide_data.get("type", "content")

        if slide_type == "hook":
            create_hook_slide(slide, slide_data, pillar)
        elif slide_type == "cta":
            create_cta_slide(slide, slide_data)
        else:
            create_content_slide(slide, slide_data, pillar)

    # Sauvegarder
    os.makedirs(PPTX_DIR, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    safe_title = "".join(c for c in title[:30] if c.isalnum() or c in "_ -")
    filename = f"carousel_{safe_title}_{timestamp}.pptx"
    filepath = os.path.join(PPTX_DIR, filename)
    prs.save(filepath)

    return filepath


# ============================================================
# TRAITEMENT DE LA QUEUE
# ============================================================
def process_carousel_queue():
    """Traite les carrousels en attente dans la queue."""
    if not os.path.exists(QUEUE_FILE):
        print("[INFO] Pas de queue trouvee")
        return []

    with open(QUEUE_FILE, "r", encoding="utf-8") as f:
        queue = json.load(f)

    generated_files = []
    modified = False

    for i, post in enumerate(queue):
        if post.get("format") == "carrousel" and post.get("status") == "ready":
            if post.get("pptx_generated"):
                continue

            print(f"\n  [CAROUSEL #{i}] Traitement...")
            content = post.get("content", "")
            pillar = post.get("pillar", "terrain")

            # Parser et structurer
            raw_slides = parse_carousel_content(content)
            if len(raw_slides) < 3:
                print(f"    [WARN] {len(raw_slides)} slides seulement (min 3) -- skip")
                continue

            structured = structure_slides(raw_slides)
            print(f"    -> {len(structured)} slides structurees")

            # Creer le PPTX
            title_slug = "".join(c for c in content[:25] if c.isalnum() or c == " ").strip()
            filepath = create_carousel_pptx(structured, title=title_slug, pillar=pillar)
            print(f"    [OK] PPTX cree: {filepath}")

            # Mettre a jour la queue
            queue[i]["pptx_generated"] = True
            queue[i]["pptx_path"] = filepath
            modified = True
            generated_files.append(filepath)

    # Sauvegarder la queue si modifiee
    if modified:
        with open(QUEUE_FILE, "w", encoding="utf-8") as f:
            json.dump(queue, f, ensure_ascii=False, indent=2)

    return generated_files


# ============================================================
# GENERATION MANUELLE (pour usage direct)
# ============================================================
def generate_from_text(text, pillar="terrain", title="manual"):
    """Genere un carrousel a partir d'un texte libre."""
    raw_slides = parse_carousel_content(text)
    if len(raw_slides) < 3:
        print(f"[ERROR] Pas assez de slides ({len(raw_slides)}). Min 3.")
        return None
    structured = structure_slides(raw_slides)
    filepath = create_carousel_pptx(structured, title=title, pillar=pillar)
    return filepath


# ============================================================
# MAIN
# ============================================================
def main():
    print(f"[START] Carousel Engine -- {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    # Traiter la queue
    print("\n[1/2] Traitement des carrousels en queue...")
    files = process_carousel_queue()
    print(f"\n       -> {len(files)} carrousels generes")

    # Lister les fichiers existants
    print(f"\n[2/2] Fichiers dans {PPTX_DIR}/:")
    if os.path.exists(PPTX_DIR):
        pptx_files = sorted(os.listdir(PPTX_DIR))
        for f in pptx_files[-10:]:
            filepath = os.path.join(PPTX_DIR, f)
            size_kb = os.path.getsize(filepath) / 1024
            print(f"       {f} ({size_kb:.0f} KB)")
        if not pptx_files:
            print("       [vide]")
    else:
        print("       [dossier inexistant]")

    print("\n[DONE] Carousel Engine termine.")


if __name__ == "__main__":
    main()

